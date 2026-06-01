#!/usr/bin/env python3
# SPDX-License-Identifier: MIT  (see LICENSE in the repo root)
"""
extract-pptx-only.py

Phase 0p extractor for the deck-import skill's PPTX-only auto-detect path.

This is the entry point used when the user supplies a .pptx without a PDF.
It walks the PPTX with the python-pptx library and emits a JSON manifest
that the auto-detect step consumes to route into the classifier and build
pipeline.

What it extracts per slide:
  - title (from the title placeholder when present)
  - body text (concatenated from text-bearing placeholders + text boxes)
  - captions (text inside picture frames, when authored)
  - speaker notes (verbatim from the notes pane)
  - embedded images, saved alongside the script with predictable names
    of the form `image-<slide_num>-<idx>.<ext>`, plus position + size in
    canvas pixels
  - slide order (raw presentation.xml order, NOT filename-sorted, because
    PowerPoint reorders slides without renumbering filenames)

Deck-level metadata:
  - title, author, subject, last-modified-by (from core properties)
  - slide_count
  - slide_size_emu and the human aspect-ratio label
  - per-slide PNG render paths (if LibreOffice is available and
    --render-pages is passed; otherwise this list is empty and the Mode A
    build step has to fall back to reconstruction or ask the user for a
    PDF render)

Output structure:
    <out_dir>/
        pptx-manifest.json       -- the JSON manifest the agent reads
        media/                   -- one file per embedded image
            image-1-1.png
            image-2-1.jpeg
            ...
        pages/                   -- (optional) per-slide PNGs from libreoffice
            page-01.png
            ...

Usage:
    python extract-pptx-only.py <pptx-path> --out /tmp/deck-import-<ts>/
    python extract-pptx-only.py deck.pptx --out /tmp/d/ --render-pages

Dependencies:
    pip install --break-system-packages python-pptx Pillow
    # Optional, for --render-pages: a working `libreoffice` on PATH.
"""

from __future__ import annotations

import argparse
import json
import os
import shutil
import subprocess
import sys
from dataclasses import dataclass, asdict
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError as e:
    sys.stderr.write(
        f"Missing dependency: {e.name}. Install with:\n"
        f"  pip install --break-system-packages python-pptx Pillow\n"
    )
    sys.exit(2)


# Canvas the deck-builder shell renders into. The build step uses these
# numbers to convert EMU positions into pixel coordinates so the same math
# matches the PDF+PPTX hybrid path.
CANVAS_W_PX = 1440
CANVAS_H_PX = 810

# Reference EMU dimensions for human aspect-ratio labels.
_REFERENCE_RATIOS = [
    ("16:9-widescreen", 12192000, 6858000),
    ("16:9-standard",   9144000,  5143500),
    ("4:3",             9144000,  6858000),
    ("16:10",           9144000,  5715000),
]


@dataclass
class ImageRecord:
    path: str
    position_px: Dict[str, float]
    size_px: Dict[str, float]
    position_emu: Dict[str, int]
    size_emu: Dict[str, int]
    content_type: str


@dataclass
class SlideRecord:
    slide_num: int
    title: str
    body_text: str
    captions: List[str]
    notes: str
    images: List[ImageRecord]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def emu_to_px(emu: int, total_emu: int, canvas_px: int) -> float:
    if not total_emu:
        return 0.0
    return round((emu / total_emu) * canvas_px, 2)


def aspect_ratio_label(cx_emu: int, cy_emu: int) -> str:
    for label, ref_cx, ref_cy in _REFERENCE_RATIOS:
        if abs(cx_emu - ref_cx) / ref_cx < 0.01 and abs(cy_emu - ref_cy) / ref_cy < 0.01:
            return label
    if cy_emu <= 0:
        return "unknown"
    return f"{round(cx_emu / cy_emu, 3)}:1 (custom)"


def safe_text(s: Optional[str]) -> str:
    return (s or "").strip()


def extension_for_image(content_type: str, blob: bytes) -> str:
    """Map a python-pptx Image content_type to a sensible file extension.
    Falls back to sniffing the first few bytes for common formats."""
    ct = (content_type or "").lower()
    if "png" in ct:
        return "png"
    if "jpeg" in ct or "jpg" in ct:
        return "jpg"
    if "gif" in ct:
        return "gif"
    if "svg" in ct:
        return "svg"
    if "tiff" in ct:
        return "tiff"
    if "bmp" in ct:
        return "bmp"
    if "webp" in ct:
        return "webp"
    # Sniff
    if blob[:8] == b"\x89PNG\r\n\x1a\n":
        return "png"
    if blob[:3] == b"\xff\xd8\xff":
        return "jpg"
    if blob[:6] in (b"GIF87a", b"GIF89a"):
        return "gif"
    if blob[:5] == b"<?xml" or blob[:4] == b"<svg":
        return "svg"
    return "bin"


# ---------------------------------------------------------------------------
# Per-shape extractors
# ---------------------------------------------------------------------------

def collect_text_from_shape(shape) -> Tuple[str, List[str]]:
    """Return (body_text_chunk, captions_chunk) for one shape.

    Title placeholders are handled separately by the caller, so this
    function returns body text for non-title text shapes plus any captions
    associated with picture frames.
    """
    body_pieces: List[str] = []
    captions: List[str] = []

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            line = "".join(run.text for run in paragraph.runs)
            line = safe_text(line)
            if line:
                body_pieces.append(line)
    # Some picture frames carry an alt-text style caption via .name or
    # description; python-pptx exposes the latter as ._element.nvPicPr/cNvPr.
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            cnv = shape._element.nvPicPr.cNvPr
            descr = cnv.get("descr")
            if descr:
                descr = safe_text(descr)
                if descr:
                    captions.append(descr)
        except (AttributeError, KeyError):
            pass

    # Recurse into groups
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            b, c = collect_text_from_shape(sub)
            if b:
                body_pieces.append(b)
            captions.extend(c)

    return "\n".join(body_pieces), captions


def collect_images_from_shape(
    shape,
    slide_num: int,
    idx_counter: List[int],
    media_dir: Path,
    sldW_emu: int,
    sldH_emu: int,
) -> List[ImageRecord]:
    """Walk one shape (and any nested group) for embedded pictures. Saves
    each to media/ with a predictable filename, returns ImageRecord entries.
    """
    out: List[ImageRecord] = []

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            image = shape.image
            blob = image.blob
            ext = extension_for_image(image.content_type, blob)
            idx_counter[0] += 1
            fname = f"image-{slide_num}-{idx_counter[0]}.{ext}"
            dest = media_dir / fname
            dest.write_bytes(blob)

            x_emu = int(shape.left or 0)
            y_emu = int(shape.top or 0)
            cx_emu = int(shape.width or 0)
            cy_emu = int(shape.height or 0)

            out.append(ImageRecord(
                path=str(dest),
                position_px={
                    "x": emu_to_px(x_emu, sldW_emu, CANVAS_W_PX),
                    "y": emu_to_px(y_emu, sldH_emu, CANVAS_H_PX),
                },
                size_px={
                    "w": emu_to_px(cx_emu, sldW_emu, CANVAS_W_PX),
                    "h": emu_to_px(cy_emu, sldH_emu, CANVAS_H_PX),
                },
                position_emu={"x": x_emu, "y": y_emu},
                size_emu={"cx": cx_emu, "cy": cy_emu},
                content_type=image.content_type or "",
            ))
        except (AttributeError, KeyError, ValueError) as exc:
            sys.stderr.write(f"WARN: failed to extract image on slide {slide_num}: {exc}\n")

    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            out.extend(collect_images_from_shape(
                sub, slide_num, idx_counter, media_dir, sldW_emu, sldH_emu
            ))

    return out


def extract_title(slide) -> str:
    """Return the slide title from the title placeholder if present."""
    try:
        title_placeholder = slide.shapes.title
    except (AttributeError, IndexError):
        return ""
    if title_placeholder is None:
        return ""
    try:
        return safe_text(title_placeholder.text)
    except AttributeError:
        return ""


def extract_notes(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    notes_slide = slide.notes_slide
    if notes_slide.notes_text_frame is None:
        return ""
    return safe_text(notes_slide.notes_text_frame.text)


# ---------------------------------------------------------------------------
# Optional: LibreOffice page rendering for Mode A page-image builds
# ---------------------------------------------------------------------------

def render_pages_with_libreoffice(pptx_path: Path, out_dir: Path) -> List[str]:
    """Convert every slide in the PPTX to a PNG using LibreOffice headless.
    Returns a list of rendered PNG paths in slide order, or [] if LibreOffice
    isn't on PATH or the conversion failed.

    Use this when the user wants a Mode A (page-image) build from PPTX-only
    input. Mode B (reconstruction) doesn't need rendered pages.
    """
    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not soffice:
        sys.stderr.write(
            "INFO: libreoffice not on PATH; skipping page rendering. "
            "Mode A builds from PPTX-only input require LibreOffice OR a PDF render.\n"
        )
        return []

    pages_dir = out_dir / "pages"
    pages_dir.mkdir(parents=True, exist_ok=True)

    try:
        # `libreoffice --convert-to png` only renders the first slide. Convert
        # to PDF first, then convert the PDF to per-page PNGs with `pdftoppm`
        # if available, falling back to ImageMagick `convert` otherwise.
        # Both fallbacks are common on Linux/macOS install bases.
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf",
             "--outdir", str(pages_dir), str(pptx_path)],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=240,
        )
        pdf_path = pages_dir / (pptx_path.stem + ".pdf")
        if not pdf_path.exists():
            sys.stderr.write("WARN: libreoffice conversion did not produce a PDF; skipping page rasters.\n")
            return []

        pdftoppm = shutil.which("pdftoppm")
        if pdftoppm:
            subprocess.run(
                [pdftoppm, "-png", "-r", "150", str(pdf_path), str(pages_dir / "page")],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=240,
            )
            # pdftoppm names files page-1.png, page-2.png, ...
            renders = sorted(pages_dir.glob("page-*.png"))
            return [str(p) for p in renders]

        convert = shutil.which("convert")
        if convert:
            subprocess.run(
                [convert, "-density", "150", str(pdf_path), str(pages_dir / "page-%02d.png")],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=240,
            )
            renders = sorted(pages_dir.glob("page-*.png"))
            return [str(p) for p in renders]

        sys.stderr.write(
            "INFO: pdftoppm and ImageMagick `convert` are both unavailable. "
            "PDF was rendered but per-page PNGs could not be produced.\n"
        )
        return []
    except subprocess.CalledProcessError as exc:
        sys.stderr.write(f"WARN: libreoffice render failed: {exc}\n")
        return []
    except subprocess.TimeoutExpired:
        sys.stderr.write("WARN: libreoffice render timed out.\n")
        return []


# ---------------------------------------------------------------------------
# Main extraction
# ---------------------------------------------------------------------------

def extract_pptx(pptx_path: Path, out_dir: Path, render_pages: bool = False) -> Dict[str, Any]:
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")
    if pptx_path.suffix.lower() != ".pptx":
        sys.stderr.write(f"WARN: expected .pptx extension, got {pptx_path.suffix}\n")

    out_dir.mkdir(parents=True, exist_ok=True)
    media_dir = out_dir / "media"
    media_dir.mkdir(parents=True, exist_ok=True)

    pres = Presentation(str(pptx_path))

    sldW_emu = int(pres.slide_width or 12192000)
    sldH_emu = int(pres.slide_height or 6858000)

    core = pres.core_properties
    metadata = {
        "title":     safe_text(core.title)    or None,
        "author":    safe_text(core.author)   or None,
        "subject":   safe_text(core.subject)  or None,
        "last_modified_by": safe_text(core.last_modified_by) or None,
        "slide_count":      len(pres.slides),
        "slide_size_emu":   {"cx": sldW_emu, "cy": sldH_emu},
        "ratio":            aspect_ratio_label(sldW_emu, sldH_emu),
    }

    slides: List[SlideRecord] = []
    for idx, slide in enumerate(pres.slides, start=1):
        title = extract_title(slide)
        body_pieces: List[str] = []
        captions: List[str] = []
        image_records: List[ImageRecord] = []
        image_idx_counter = [0]

        for shape in slide.shapes:
            # Title placeholder is already captured separately, skip its body.
            try:
                is_title = (shape == slide.shapes.title)
            except (AttributeError, IndexError):
                is_title = False
            if not is_title:
                body, caps = collect_text_from_shape(shape)
                if body:
                    body_pieces.append(body)
                captions.extend(caps)

            image_records.extend(collect_images_from_shape(
                shape, idx, image_idx_counter, media_dir, sldW_emu, sldH_emu,
            ))

        slides.append(SlideRecord(
            slide_num=idx,
            title=title,
            body_text="\n".join(b for b in body_pieces if b),
            captions=captions,
            notes=extract_notes(slide),
            images=image_records,
        ))

    rendered_pages: List[str] = []
    if render_pages:
        rendered_pages = render_pages_with_libreoffice(pptx_path, out_dir)

    manifest: Dict[str, Any] = {
        "source_pptx":  str(pptx_path.resolve()),
        "extracted_at": datetime.now(timezone.utc).isoformat(),
        "metadata":     metadata,
        "slides":       [_slide_to_dict(s) for s in slides],
        "rendered_pages": rendered_pages,
    }

    manifest_path = out_dir / "pptx-manifest.json"
    manifest_path.write_text(json.dumps(manifest, indent=2))
    return manifest


def _slide_to_dict(s: SlideRecord) -> Dict[str, Any]:
    d = asdict(s)
    # Convert ImageRecord dataclasses back to plain dicts (asdict already does
    # this; this stub stays in case we add transforms later).
    return d


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main() -> int:
    parser = argparse.ArgumentParser(
        description="Extract text + images + notes from a .pptx into a JSON manifest."
    )
    parser.add_argument("pptx_path", type=Path, help="Path to the source .pptx file.")
    parser.add_argument("--out", type=Path, required=True,
                        help="Output directory (will be created).")
    parser.add_argument("--render-pages", action="store_true",
                        help="Also render every slide to PNG via LibreOffice "
                             "(needed for Mode A page-image builds from "
                             "PPTX-only input).")
    args = parser.parse_args()

    try:
        manifest = extract_pptx(args.pptx_path, args.out, render_pages=args.render_pages)
    except Exception as exc:  # noqa: BLE001
        sys.stderr.write(f"Extraction failed: {exc}\n")
        return 1

    summary = {
        "slide_count":      manifest["metadata"]["slide_count"],
        "ratio":            manifest["metadata"]["ratio"],
        "manifest":         str(args.out / "pptx-manifest.json"),
        "media_dir":        str(args.out / "media"),
        "rendered_pages":   len(manifest["rendered_pages"]),
    }
    print(json.dumps(summary, indent=2))
    return 0


if __name__ == "__main__":
    sys.exit(main())
